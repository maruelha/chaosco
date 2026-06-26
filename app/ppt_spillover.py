"""Spillover Status Report — PowerPoint builder."""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from app.ppt_utils import (
    _i, _rgb, _add_rect, _add_text,
    INK, SUB, FAINT, LINE, HAIR, CARDBG, WHITE,
    GREEN, RED, ORANGE, BLUE, PURPLE,
    HEAD, BODY, PAGEW, PAGEH, M, CW,
)

# ---------------------------------------------------------------------------
# Spillover-specific theme — change here to experiment without touching retail
# ---------------------------------------------------------------------------

TITLE = "SPILLOVER STATUS / DTC"
TAG   = "UAT COORDINATION · CORE SOUTH"

CONTENT_Y  = 1.22   # content starts below header
SECTION_H  = 0.40   # height of a section banner
ITEM_H     = 1.08   # height of each item row
BOTTOM_PAD = 0.28   # reserved space at bottom of slide

# Item row column layout
_NUM_W    = 0.45
_BADGE_W  = 1.80
_ORDERS_W = 2.50
_GAP_NUM  = 0.05    # gap: num → badge
_GAP_BADGE = 0.20   # gap: badge → text
_GAP_TEXT  = 0.10   # gap: text → orders
_TEXT_W   = CW - _NUM_W - _GAP_NUM - _BADGE_W - _GAP_BADGE - _ORDERS_W - _GAP_TEXT
_X_NUM    = M
_X_BADGE  = M + _NUM_W + _GAP_NUM
_X_TEXT   = _X_BADGE + _BADGE_W + _GAP_BADGE
_X_ORDERS = M + CW - _ORDERS_W

# Stat chips (top-right header area, 4 chips in a row)
_CHIP_W = 1.10
_CHIP_H = 0.72
_CHIP_G = 0.12
_CHIPS_W = 4 * _CHIP_W + 3 * _CHIP_G
_CHIP_X0 = M + CW - _CHIPS_W  # left edge of first chip
_CHIP_Y  = 0.18

# ---------------------------------------------------------------------------
# Status badge mapping  (key = status field value, lower-cased)
# unknown statuses fall back to a grey badge showing the raw text
# ---------------------------------------------------------------------------

_STATUS_STYLE: dict[str, tuple] = {
    # key                      bg colour          icon   display label
    "blocked dtc":         (RED,                 "🚫",  "BLOCKED DTC"),
    "clarification":       (PURPLE,              "💬",  "CLARIFICATION"),
    "to sales":            (GREEN,               "↩️",  "TO SALES"),
    "sales in progress":   (ORANGE,              "🏃",  "SALES IN PROGRESS"),
    "waiting sf creation": (_rgb("6E6E73"),      "⏰",  "WAITING SF CREATION"),
    "passed":              (_rgb("248A3D"),       "✅",  "PASSED"),
    "in progress":         (BLUE,                "⚙️",  "IN PROGRESS"),
    "open":                (_rgb("1D1D1F"),       "📋",  "OPEN"),
}

# ---------------------------------------------------------------------------
# Section definitions — order drives slide order
# ---------------------------------------------------------------------------

_SECTIONS = [
    {"key": "yes",      "label": "CRITICAL FOR SIGN-OFF", "icon": "⚠️",  "bg": RED,              "fg": WHITE},
    {"key": "slightly", "label": "SLIGHTLY CRITICAL",      "icon": "⚡",  "bg": ORANGE,           "fg": WHITE},
    {"key": "no",       "label": "NON-CRITICAL",           "icon": "↩️", "bg": _rgb("248A3D"),   "fg": WHITE},
    {"key": "unset",    "label": "NOT CATEGORISED",        "icon": "❓", "bg": _rgb("6E6E73"),   "fg": WHITE},
]


def _crit_key(val: str | None) -> str:
    if not val:
        return "unset"
    v = val.strip().lower()
    if v == "yes":
        return "yes"
    if v == "slightly":
        return "slightly"
    if v == "no":
        return "no"
    return "unset"


# ---------------------------------------------------------------------------
# Header
# ---------------------------------------------------------------------------

def _draw_header(slide, today: str, total: int, n_crit: int,
                 n_slightly: int, n_non_crit: int, continued: bool = False):
    # Small tag
    _add_text(slide, TAG, M, 0.15, 7.0, 0.22,
              font=BODY, size=8, bold=True, color=FAINT)

    # Large title
    _add_text(slide, f"🔥  {TITLE}", M, 0.33, 7.0, 0.52,
              font=HEAD, size=24, bold=True, color=INK)

    # Subtitle
    sub = f"⏰  End of day  ·  {today}" + ("  ·  continued" if continued else "")
    _add_text(slide, sub, M, 0.86, 7.0, 0.24,
              font=BODY, size=10, color=SUB)

    # Stat chips (top-right)
    chips = [
        {"label": "OPEN",         "n": total,      "color": INK},
        {"label": "CRITICAL",     "n": n_crit,     "color": RED},
        {"label": "SLIGHTLY",     "n": n_slightly, "color": ORANGE},
        {"label": "NON-CRITICAL", "n": n_non_crit, "color": _rgb("6E6E73")},
    ]
    for i, chip in enumerate(chips):
        x = _CHIP_X0 + i * (_CHIP_W + _CHIP_G)
        _add_rect(slide, x, _CHIP_Y, _CHIP_W, _CHIP_H,
                  fill=CARDBG, line=HAIR, rounded=True)
        _add_text(slide, str(chip["n"]),
                  x, _CHIP_Y + 0.08, _CHIP_W, 0.32,
                  font=HEAD, size=18, bold=True, color=chip["color"],
                  align=PP_ALIGN.CENTER)
        _add_text(slide, chip["label"],
                  x, _CHIP_Y + 0.44, _CHIP_W, 0.22,
                  font=BODY, size=7, bold=True, color=SUB,
                  align=PP_ALIGN.CENTER)

    # Divider
    _add_rect(slide, M, 1.12, CW, 0.01, fill=LINE)


# ---------------------------------------------------------------------------
# Section banner
# ---------------------------------------------------------------------------

def _draw_section_banner(slide, y: float, section: dict):
    _add_rect(slide, M, y, CW, SECTION_H, fill=section["bg"])
    _add_text(slide, f"{section['icon']}  {section['label']}",
              M + 0.22, y + 0.10, CW - 0.44, SECTION_H - 0.16,
              font=BODY, size=11, bold=True, color=section["fg"])


# ---------------------------------------------------------------------------
# Item row
# ---------------------------------------------------------------------------

def _draw_item_row(slide, y: float, row_num: int, item: dict,
                   orders: list[dict], is_alt: bool):
    row_bg = _rgb("F8F8F8") if is_alt else WHITE
    _add_rect(slide, M, y, CW, ITEM_H, fill=row_bg, line=HAIR)

    # Row number
    _add_text(slide, f"{row_num:02d}",
              _X_NUM, y + (ITEM_H - 0.34) / 2, _NUM_W, 0.34,
              font=HEAD, size=13, bold=True, color=FAINT,
              align=PP_ALIGN.CENTER)

    # Status badge
    status_raw = (item.get("status") or "").strip()
    style      = _STATUS_STYLE.get(status_raw.lower())
    if style:
        bg_col, icon, label = style
    else:
        bg_col = _rgb("AEAEB2")
        icon   = "•"
        label  = status_raw.upper()[:20] if status_raw else "—"

    badge_y = y + (ITEM_H - 0.38) / 2
    _add_rect(slide, _X_BADGE, badge_y, _BADGE_W, 0.38, fill=bg_col, rounded=True)
    _add_text(slide, f"{icon}  {label}",
              _X_BADGE + 0.10, badge_y + 0.07, _BADGE_W - 0.16, 0.26,
              font=BODY, size=8, bold=True, color=WHITE)

    # Name
    _add_text(slide, item.get("name") or "—",
              _X_TEXT, y + 0.08, _TEXT_W, 0.36,
              font=BODY, size=10, bold=True, color=INK)

    # Next step
    next_step = (item.get("next_step") or "").strip()
    if next_step:
        _add_text(slide, next_step,
                  _X_TEXT, y + 0.46, _TEXT_W, 0.54,
                  font=BODY, size=9, italic=True, color=SUB)

    # Order numbers (right column)
    if orders:
        lines = []
        for od in orders[:4]:
            num = (od.get("order_number") or "").strip()
            typ = (od.get("order_type") or "").strip()
            s4  = " ✓" if od.get("docs_in_s4") else ""
            if num:
                lines.append(f"{typ + ':  ' if typ else ''}{num}{s4}")
        if lines:
            _add_text(slide, "\n".join(lines),
                      _X_ORDERS, y + 0.12, _ORDERS_W, ITEM_H - 0.24,
                      font=BODY, size=8, color=SUB)


# ---------------------------------------------------------------------------
# Pagination
# ---------------------------------------------------------------------------

def _paginate(items: list[dict]) -> list[list[dict]]:
    """Return a list of slides; each slide is a list of element dicts."""
    groups: dict[str, list[dict]] = {}
    for item in items:
        k = _crit_key(item.get("critical_for_signoff"))
        groups.setdefault(k, []).append(item)

    # Build flat element list in section order
    elements: list[dict] = []
    row_num = 0
    for sec in _SECTIONS:
        key = sec["key"]
        if key not in groups:
            continue
        elements.append({"type": "section", "section": sec})
        for item in groups[key]:
            row_num += 1
            elements.append({"type": "item", "item": item, "row_num": row_num})

    # Fill slides greedily; a section banner must be followed by at least 1 item
    slides: list[list[dict]] = []
    current: list[dict] = []
    remaining = PAGEH - CONTENT_Y - BOTTOM_PAD

    for elem in elements:
        if elem["type"] == "section":
            needed = SECTION_H + ITEM_H  # don't orphan a banner at bottom
            if current and remaining < needed:
                slides.append(current)
                current = []
                remaining = PAGEH - CONTENT_Y - BOTTOM_PAD
            current.append(elem)
            remaining -= SECTION_H
        else:
            if remaining < ITEM_H:
                slides.append(current)
                current = []
                remaining = PAGEH - CONTENT_Y - BOTTOM_PAD
            current.append(elem)
            remaining -= ITEM_H

    if current:
        slides.append(current)

    return slides


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def build_spillover_ppt(items: list[dict],
                        order_details: dict[int, list[dict]],
                        today: str) -> bytes:
    """Build the spillover status PPT and return raw bytes.

    items         — list of dicts from database.get_spillover_report_items(),
                    sorted by criticality (web.py does this before calling)
    order_details — {spillover_id: [order_detail dicts]}
    today         — date string (YYYY-MM-DD)
    """
    n_crit     = sum(1 for i in items if _crit_key(i.get("critical_for_signoff")) == "yes")
    n_slightly = sum(1 for i in items if _crit_key(i.get("critical_for_signoff")) == "slightly")
    n_non_crit = sum(1 for i in items if _crit_key(i.get("critical_for_signoff")) == "no")
    total      = len(items)

    slide_pages = _paginate(items)
    n_pages     = len(slide_pages)

    prs = Presentation()
    prs.slide_width  = int(Inches(PAGEW))
    prs.slide_height = int(Inches(PAGEH))

    item_counter = 0  # track alternating row colour across section banners

    for pi, page_elems in enumerate(slide_pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE

        _draw_header(s, today, total, n_crit, n_slightly, n_non_crit,
                     continued=(pi > 0))

        # Page number (bottom-right)
        _add_text(s, f"PAGE {pi + 1} / {n_pages}",
                  M, PAGEH - 0.24, CW, 0.20,
                  font=BODY, size=8, color=FAINT, align=PP_ALIGN.RIGHT)

        y = CONTENT_Y
        for elem in page_elems:
            if elem["type"] == "section":
                _draw_section_banner(s, y, elem["section"])
                y += SECTION_H
            else:
                item    = elem["item"]
                sid     = item.get("spillover_id")
                orders  = order_details.get(sid, [])
                is_alt  = (item_counter % 2 == 1)
                _draw_item_row(s, y, elem["row_num"], item, orders, is_alt)
                item_counter += 1
                y += ITEM_H

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
