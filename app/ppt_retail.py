"""Retail Status Report — PowerPoint builder."""
from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.ppt_utils import (
    _i, _add_rect, _add_text, _add_header,
    INK, SUB, FAINT, LINE, HAIR, CARDBG, WHITE, INK3A, ROWALT,
    BLUE, ORANGE, GREEN, GREENL, RED, _rgb,
    HEAD, BODY, PAGEW, PAGEH, M, CW,
)

ROWS_PER_SLIDE = 12
DASH = "—"

BREAKDOWN_CARDS = [
    {"label": "Incoming (Gatekeeper)", "owner": "Jose",      "icon": "📥", "color": _rgb("0071E3"), "tint": _rgb("E8F2FD")},
    {"label": "Ready for validation",  "owner": "Key users", "icon": "🎯", "color": _rgb("248A3D"), "tint": _rgb("E9F6EE")},
    {"label": "In Progress",           "owner": "Key users", "icon": "⚙️", "color": _rgb("B36200"), "tint": _rgb("FBF0E2")},
    {"label": "In Clarification",      "owner": "Key users", "icon": "💬", "color": _rgb("6E56CF"), "tint": _rgb("EFEBFB")},
    {"label": "Blocked",               "owner": "Tech team", "icon": "🚫", "color": _rgb("C0392B"), "tint": _rgb("FBEBEA")},
]

BUCKET_KEY_TO_CARD = {
    "incoming_gatekeeper":  0,
    "ready_for_validation": 1,
    "in_progress":          2,
    "in_clarification":     3,
    "blocked":              4,
}

TITLE = "Retail Status Report — MB-DTC Viewpoint"


# ---------------------------------------------------------------------------
# Slide 1 — Summary
# ---------------------------------------------------------------------------

def _slide1_summary(prs, report: dict, today: str, total_test_cases: int,
                    missing_categories: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE

    b   = report["buckets"]
    pct = round((b["passed_with_dtc"] / total_test_cases) * 100, 1)

    _add_header(s, TITLE,
                f"Computed live · {today} · {b['with_dtc'] + b['back_with_sales']} of {total_test_cases} tracked")

    # ── Stats strip ──────────────────────────────────────────────────────────
    strip_y, strip_h = 1.28, 1.0
    _add_rect(s, M, strip_y, CW, strip_h, fill=_rgb("F5F5F7"), rounded=True)

    sw = 1.9
    def stat(x, num, label, color):
        _add_text(s, str(num), x, strip_y + 0.12, sw, 0.42,
                  font=HEAD, size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_text(s, label.upper(), x, strip_y + 0.57, sw, 0.32,
                  font=BODY, size=11, bold=True, color=SUB, align=PP_ALIGN.CENTER)

    stat(M + 0.2,        total_test_cases, "Total test cases", INK)
    stat(M + 0.2 + sw,   b["with_dtc"] + b["back_with_sales"], "In tracker", INK)
    stat(M + 0.2 + sw*2, f"{pct}%",        "Passed of total", GREEN)

    _add_rect(s, M + 0.3 + sw*3, strip_y + 0.18, 0.01, strip_h - 0.36, fill=LINE)
    note_x = M + 0.55 + sw * 3
    _add_text(s, "Additional test cases on top of total test cases:",
              note_x, strip_y + 0.12, CW - (note_x - M) - 0.2, 0.28,
              font=BODY, size=9, bold=True, color=INK3A, align=PP_ALIGN.LEFT)
    bullets = "\n".join(f"• {c}" for c in missing_categories)
    _add_text(s, bullets,
              note_x + 0.05, strip_y + 0.42, CW - (note_x - M) - 0.25, 0.48,
              font=BODY, size=9, color=SUB, align=PP_ALIGN.LEFT)

    # ── Overview buckets ─────────────────────────────────────────────────────
    y = strip_y + strip_h + 0.24
    _add_text(s, "Overview", M, y, CW, 0.32, font=HEAD, size=15, bold=True, color=INK)
    y += 0.40

    buckets = [
        {"n": b["back_with_sales"],      "label": "Back with Sales",      "color": INK,    "border": LINE,              "sub": None},
        {"n": b["with_dtc"],             "label": "With DTC",             "color": BLUE,   "border": BLUE,              "sub": f"{b['in_progress_with_dtc']} in progress  ·  {b['passed_with_dtc']} passed"},
        {"n": b["in_progress_with_dtc"], "label": "In Progress with DTC", "color": ORANGE, "border": _rgb("FF9500"),    "sub": None},
        {"n": b["passed_with_dtc"],      "label": "Passed with DTC",      "color": GREEN,  "border": GREENL,            "sub": None},
    ]
    gap = 0.25; bw = (CW - gap * 3) / 4; bh = 1.5
    for i, bkt in enumerate(buckets):
        x = M + i * (bw + gap)
        _add_rect(s, x, y, bw, bh, fill=WHITE, line=bkt["border"], rounded=True)
        _add_text(s, str(bkt["n"]), x, y + 0.15, bw, 0.62,
                  font=HEAD, size=32, bold=True, color=bkt["color"], align=PP_ALIGN.CENTER)
        _add_text(s, bkt["label"].upper(), x + 0.1, y + 0.82, bw - 0.2, 0.34,
                  font=BODY, size=11, bold=True, color=INK3A, align=PP_ALIGN.CENTER)
        if bkt["sub"]:
            _add_text(s, bkt["sub"], x + 0.1, y + 1.18, bw - 0.2, 0.24,
                      font=BODY, size=9, color=SUB, align=PP_ALIGN.CENTER)
    y += bh + 0.14

    # ── In-Progress Breakdown ─────────────────────────────────────────────────
    _add_text(s, "In-Progress Breakdown", M, y, CW, 0.30,
              font=HEAD, size=15, bold=True, color=INK)
    y += 0.36

    cgap = 0.22; cw = (CW - cgap * 4) / 5; ch = 1.15
    breakdown_vals = [
        b["incoming_gatekeeper"], b["ready_for_validation"],
        b["in_progress"], b["in_clarification"], b["blocked"],
    ]
    for i, (card, val) in enumerate(zip(BREAKDOWN_CARDS, breakdown_vals)):
        x = M + i * (cw + cgap)
        _add_rect(s, x, y, cw, ch, fill=CARDBG, line=HAIR, rounded=True)
        _add_rect(s, x + cw/2 - 0.23, y + 0.10, 0.46, 0.46,
                  fill=card["tint"], rounded=False)
        _add_text(s, card["icon"], x + cw/2 - 0.18, y + 0.12, 0.36, 0.32,
                  font=BODY, size=14, color=card["color"], align=PP_ALIGN.CENTER)
        num_color = card["color"] if val > 0 else SUB
        _add_text(s, str(val), x, y + 0.54, cw, 0.30,
                  font=HEAD, size=18, bold=True, color=num_color, align=PP_ALIGN.CENTER)
        _add_text(s, card["label"], x + 0.06, y + 0.82, cw - 0.12, 0.22,
                  font=BODY, size=10, bold=True, color=INK3A, align=PP_ALIGN.CENTER)
        _add_text(s, card["owner"], x + 0.06, y + 1.01, cw - 0.12, 0.16,
                  font=BODY, size=9, color=SUB, align=PP_ALIGN.CENTER)
    y += ch + 0.16

    # ── Comments box ─────────────────────────────────────────────────────────
    cmt_h = PAGEH - 0.45 - y
    if cmt_h > 0.2:
        _add_rect(s, M, y, CW, cmt_h, fill=WHITE, line=LINE, rounded=True)
        _add_text(s, "Comments", M + 0.15, y + 0.08, 2.0, 0.22,
                  font=BODY, size=8, bold=True, color=FAINT, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide 2+ — Defects table
# ---------------------------------------------------------------------------

def _slide_defects(prs, defects: list[dict], today: str,
                   mb_total: int, sales_total: int, impacted_total: int):
    col_w = [0.45, 1.25, 2.75, 2.30, 1.30, 1.65, 1.00, 1.40]
    col_x = []
    cx = M
    for w in col_w:
        col_x.append(cx)
        cx += w
    headers = ["#", "DEFECT ID", "SOLMAN NAME", "ASSIGNED TO",
               "DATE REPORTED", "STATUS", "IMPACTED MB", "IMPACTED SALES"]

    chunks = [defects[i:i+ROWS_PER_SLIDE] for i in range(0, max(len(defects), 1), ROWS_PER_SLIDE)]

    for ci, chunk in enumerate(chunks):
        is_last = ci == len(chunks) - 1
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = WHITE

        title_suffix = f"  ({ci+1}/{len(chunks)})" if len(chunks) > 1 else ""
        _add_header(s, TITLE,
                    f"Retail Defects — Active (excl. Confirmed & Withdrawn) · {today}")
        _add_text(s, "Retail Defects — Active (excl. Confirmed & Withdrawn)" + title_suffix,
                  M, 1.16, CW, 0.30, font=HEAD, size=15, bold=True, color=INK)
        _add_text(s, f"{len(defects)} active defects  ·  Impacted = test cases referencing the defect and not yet passed  ·  MB = DTC O2C follow-up",
                  M, 1.48, CW, 0.24, font=BODY, size=9, italic=True, color=SUB)

        ty = 1.78; hH = 0.34
        for i, h in enumerate(headers):
            num_col  = i >= 6
            idx_col  = i == 0
            align = PP_ALIGN.RIGHT if num_col else (PP_ALIGN.CENTER if idx_col else PP_ALIGN.LEFT)
            _add_text(s, h, col_x[i], ty, col_w[i], hH,
                      font=BODY, size=8, bold=True, color=SUB, align=align)
        _add_rect(s, M, ty + hH, CW, 0.01, fill=LINE)

        row_h = 0.32; ty += hH + 0.04
        for ri, d in enumerate(chunk):
            ry = ty + ri * row_h
            row_num = ci * ROWS_PER_SLIDE + ri + 1
            muted  = d.get("impacted_tc_count", 0) == 0
            tcol   = FAINT if muted else INK
            scol   = FAINT if muted else SUB

            if ri % 2 == 1:
                _add_rect(s, M, ry, CW, row_h, fill=ROWALT)

            _add_text(s, str(row_num),                    col_x[0], ry, col_w[0], row_h, font=BODY, size=9, color=scol,  align=PP_ALIGN.CENTER)
            _add_text(s, d.get("defect_id") or "",        col_x[1], ry, col_w[1], row_h, font=BODY, size=9, color=BLUE if not muted else FAINT, align=PP_ALIGN.LEFT)
            _add_text(s, d.get("solman_name") or DASH,    col_x[2], ry, col_w[2], row_h, font=BODY, size=9, color=tcol, align=PP_ALIGN.LEFT)
            _add_text(s, d.get("assigned_to") or DASH,    col_x[3], ry, col_w[3], row_h, font=BODY, size=9, color=scol, align=PP_ALIGN.LEFT)
            _add_text(s, d.get("date_reported") or DASH,  col_x[4], ry, col_w[4], row_h, font=BODY, size=9, color=scol, align=PP_ALIGN.LEFT)
            _add_text(s, d.get("solman_status") or DASH,  col_x[5], ry, col_w[5], row_h, font=BODY, size=9, color=scol, align=PP_ALIGN.LEFT)

            mb    = d.get("impacted_tc_count") if d.get("dtco2c")     else None
            sales = d.get("impacted_tc_count") if not d.get("dtco2c") else None
            _add_text(s, str(mb)    if mb    else DASH, col_x[6], ry, col_w[6], row_h, font=BODY, size=9, bold=bool(mb),    color=tcol if mb    else FAINT, align=PP_ALIGN.RIGHT)
            _add_text(s, str(sales) if sales else DASH, col_x[7], ry, col_w[7], row_h, font=BODY, size=9, bold=bool(sales), color=tcol if sales else FAINT, align=PP_ALIGN.RIGHT)

        after = ty + len(chunk) * row_h + 0.05

        if is_last:
            _add_rect(s, M, after, CW, 0.01, fill=LINE)
            after += 0.08
            _add_text(s, "Total impacted test cases", col_x[5], after, col_w[5], 0.30,
                      font=BODY, size=9, bold=True, color=INK, align=PP_ALIGN.LEFT)
            _add_text(s, str(mb_total), col_x[6], after, col_w[6], 0.30,
                      font=BODY, size=9, bold=True, color=INK, align=PP_ALIGN.RIGHT)
            _add_text(s, str(sales_total),  col_x[7], after, col_w[7], 0.30,
                      font=BODY, size=9, bold=True, color=INK, align=PP_ALIGN.RIGHT)
            after += 0.32
            _add_text(s, "Combined total", col_x[5], after, col_w[5], 0.30,
                      font=BODY, size=9, color=SUB, align=PP_ALIGN.LEFT)
            _add_text(s, str(impacted_total), col_x[7], after, col_w[7], 0.30,
                      font=BODY, size=9, bold=True, color=INK, align=PP_ALIGN.RIGHT)
            after += 0.28

            sg = 0.30; scw = (CW - sg * 2) / 3; sch = 0.90
            defect_count = sum(1 for d in defects if d.get("impacted_tc_count", 0) > 0)
            sum_cards = [
                {"n": mb_total,       "label": "Impacted Test Cases — MB",    "sub": "Our follow-up",       "sub_color": ORANGE},
                {"n": sales_total,    "label": "Impacted Test Cases — Sales", "sub": "Sales follow-up",     "sub_color": ORANGE},
                {"n": impacted_total, "label": "Total Impacted Test Cases",   "sub": f"from {defect_count} defects", "sub_color": BLUE},
            ]
            for i, c in enumerate(sum_cards):
                x = M + i * (scw + sg)
                _add_rect(s, x, after, scw, sch, fill=CARDBG, line=HAIR, rounded=True)
                _add_text(s, str(c["n"]), x, after + 0.08, scw, 0.38,
                          font=HEAD, size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)
                _add_text(s, c["label"], x + 0.1, after + 0.46, scw - 0.2, 0.28,
                          font=BODY, size=11, bold=True, color=INK3A, align=PP_ALIGN.CENTER)
                _add_text(s, c["sub"],   x + 0.1, after + 0.72, scw - 0.2, 0.20,
                          font=BODY, size=10, color=c["sub_color"], align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def build_retail_ppt(report: dict, impacted_defects: list[dict],
                     mb_total: int, sales_total: int, impacted_total: int,
                     total_test_cases: int, today: str,
                     missing_categories: list[str]) -> bytes:
    """Build the retail status PPT and return raw bytes."""
    prs = Presentation()
    prs.slide_width  = int(Inches(PAGEW))
    prs.slide_height = int(Inches(PAGEH))

    _slide1_summary(prs, report, today, total_test_cases, missing_categories)
    _slide_defects(prs, impacted_defects, today, mb_total, sales_total, impacted_total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
