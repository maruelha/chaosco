"""Shared primitives for all PPT report builders.

Each report builder (ppt_retail.py, ppt_spillover.py, …) imports what it
needs from here and can override any constant locally for a different look.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Colour helpers
# ---------------------------------------------------------------------------

def _rgb(hex6: str) -> RGBColor:
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Default palette — override in individual builders if needed
INK    = _rgb("1D1D1F")
SUB    = _rgb("6E6E73")
FAINT  = _rgb("AEAEB2")
LINE   = _rgb("D2D2D7")
HAIR   = _rgb("E5E5EA")
CARDBG = _rgb("FAFAFA")
STRIP  = _rgb("F5F5F7")
ROWALT = _rgb("FAFAFA")
WHITE  = _rgb("FFFFFF")
INK3A  = _rgb("3A3A3C")

BLUE   = _rgb("0071E3")
ORANGE = _rgb("B36200")
GREEN  = _rgb("248A3D")
GREENL = _rgb("34C759")
RED    = _rgb("C0392B")
PURPLE = _rgb("6E56CF")

# ---------------------------------------------------------------------------
# Typography
# ---------------------------------------------------------------------------

HEAD = "Bookman Old Style"
BODY = "Calibri"

# ---------------------------------------------------------------------------
# Slide dimensions (inches) — all reports use the same canvas
# ---------------------------------------------------------------------------

PAGEW = 13.3
PAGEH = 7.5
M     = 0.6          # left/right margin
CW    = PAGEW - M * 2  # content width

# ---------------------------------------------------------------------------
# Low-level drawing primitives
# ---------------------------------------------------------------------------

def _i(v: float) -> int:
    """Inches to EMU."""
    return int(Inches(v))


def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _add_rect(slide, x, y, w, h, fill: RGBColor | None = None,
              line: RGBColor | None = None, line_pt: float = 1.0,
              rounded: bool = False):
    if rounded:
        sp = slide.shapes.add_shape(
            9,  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            _i(x), _i(y), _i(w), _i(h)
        )
        try:
            sp.adjustments[0] = 0.04
        except (IndexError, Exception):
            pass
    else:
        sp = slide.shapes.add_shape(
            1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
            _i(x), _i(y), _i(w), _i(h)
        )

    if fill:
        _solid(sp, fill)
    else:
        sp.fill.background()

    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_pt)
    else:
        _no_line(sp)

    return sp


def _add_text(slide, text: str, x, y, w, h,
              font=BODY, size=10, bold=False, italic=False, color=INK,
              align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(_i(x), _i(y), _i(w), _i(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    f = run.font
    f.name   = font
    f.size   = Pt(size)
    f.bold   = bold
    f.italic = italic
    f.color.rgb = color
    return txb


def _add_header(slide, title: str, sub: str):
    """Standard report header: large title + muted subtitle + divider line."""
    _add_text(slide, title,
              M, 0.32, CW, 0.5, font=HEAD, size=24, bold=True, color=INK,
              align=PP_ALIGN.LEFT)
    _add_text(slide, sub,
              M, 0.80, CW, 0.28, font=BODY, size=11, color=FAINT,
              align=PP_ALIGN.LEFT)
    _add_rect(slide, M, 1.12, CW, 0.01, fill=LINE)
